import streamlit as st
import os


from pptx import Presentation
from huggingface_hub import InferenceClient
client = InferenceClient(
    "mistralai/Mixtral-8x7B-Instruct-v0.1"
)

temperature=0.7
max_new_tokens=2000
top_p=0.95
repetition_penalty=1.0

temperature = float(temperature)
if temperature < 1e-2:
    temperature = 1e-2
top_p = float(top_p)

generate_kwargs = dict(
    temperature=temperature,
    max_new_tokens=max_new_tokens,
    top_p=top_p,
    repetition_penalty=repetition_penalty,
    do_sample=True,
    seed=42,
)

def extract_text_from_slide(slide):
    text = ''
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            text += shape.text + '\n'
        elif shape.shape_type == 13:  # Check if shape is an image (13 is the shape ID for images)
            if hasattr(shape, 'alt_text') and shape.alt_text:
                text += shape.alt_text + '\n'
    return text.strip()
def extract_contents_from_ppt(ppt_file):
    presentation = Presentation(ppt_file)
    slides_contents = {}
    for i, slide in enumerate(presentation.slides, start=1):
        slides_contents[f"Slide {i}"] = extract_text_from_slide(slide)
    return slides_contents
def concatenate_dictionary_contents(dictionary):
    concatenated_string = ""
    for key, value in dictionary.items():
        concatenated_string += f"{key}:\n{value}\n\n"
    return concatenated_string
def generate_response(prompt):
    stream = client.text_generation(prompt, **generate_kwargs, return_full_text=False)
    return stream


def main():
    st.set_page_config(page_title="Introvert", layout="wide")
    
    st.title("Get a script for your presentation!")
    
    uploaded_file = st.file_uploader("Upload a PowerPoint file", type=["pptx"])
    contents = extract_contents_from_ppt(uploaded_file)
    user_message = st.chat_message("user", avatar='🦖')
    if uploaded_file is not None:

        with st.spinner("Generating script..."):
            prompt=f"Given the slide-wise contents of my presentation, create a script that maintains engagement and clarity while delivering a 15-20 minute presentation to higher-ups.  DONOT put the references slide in the script. Add some points of your own too related to the per slide contents The required context is within backticks, use only that to create the script. DO NOT generate a script that sounds robotic or overly formal. Use a conversational tone and simple language.The script should offer explanations and insights beyond the slide contents. PPT CONTENTS ```{concatenate_dictionary_contents(contents)}```  Give additional script explaining each slide in depth, about the content. Give MINIMUM 3-4 sentences per slide. *MAKE SURE TO USE EVERY SINGLE PROVIDED CONTEXT DETAIL.* "
            
            response=generate_response(prompt)
            
        user_message.write(response)

if __name__ == "__main__":
    main()
# Given the slide wise contents of my ppt, give me a script i can follow so that i dont have unnecessary breaks in my flow while presenting to the higher ups. The provided script should be about for 15-20 minutes of presentation.DONOT put the references slide in the script. MAKE SURE TO USE EVERY SINGLE PROVIDED CONTEXT DETAIL. Give minimum 3-4 sentences per slide. Add some points of your own too related to the per slide contents  The required context is within backticks, use only that to create the script.  DO NOT generate a script which sounds non human. Use easy words and a humanly written tone. PPT CONTENTS 